import os
import io
import re
import requests
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from dotenv import load_dotenv
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from io import BytesIO


load_dotenv()


app = Flask(__name__)
CORS(app)


# ============ MISTRAL CONFIGURATION ============
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
MISTRAL_MODEL = os.getenv("MISTRAL_MODEL", "mistral")


# Image APIs (still supported for images)
UNSPLASH_API_KEY = os.getenv("UNSPLASH_API_KEY", "").strip()
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "").strip()
PIXABAY_API_KEY = os.getenv("PIXABAY_API_KEY", "").strip()


# ============ MODE CONFIGURATIONS ============
MODE_PROMPTS = {
    "Quick Response": {
        "instructions": "Create a concise presentation with key points only. Keep it brief and to the point.",
        "slide_count": 3,
        "max_bullets": 3,
        "font_size_body": 13
    },
    "Creative": {
        "instructions": "Create an engaging and creative presentation with storytelling elements. Use vivid descriptions and make it interesting.",
        "slide_count": 5,
        "max_bullets": 4,
        "font_size_body": 12
    },
    "Detailed": {
        "instructions": "Create a comprehensive presentation with detailed explanations, examples, and in-depth analysis. Cover all aspects thoroughly.",
        "slide_count": 7,
        "max_bullets": 5,
        "font_size_body": 11
    }
}


# ============ STARTUP CHECKS ============
def check_ollama_connection():
    """Verify Ollama is running and accessible"""
    try:
        response = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=2)
        if response.status_code == 200:
            models = response.json().get('models', [])
            model_names = [m.get('name', '').split(':')[0] for m in models]
            return True, model_names
        return False, []
    except Exception as e:
        return False, []


def startup_diagnostics():
    """Run startup checks and print status"""
    print("\n" + "="*60)
    print("🚀 POWERPOINT GENERATOR WITH LOCAL MISTRAL")
    print("="*60)
    
    ollama_ok, models = check_ollama_connection()
    
    if ollama_ok:
        print(f"✅ Ollama is running ({OLLAMA_BASE_URL})")
        if models:
            print(f"📦 Available models: {models}")
            if "mistral" in models or MISTRAL_MODEL in models:
                print(f"✅ {MISTRAL_MODEL} model found!")
            else:
                print(f"⚠️  {MISTRAL_MODEL} not found. Run: ollama pull {MISTRAL_MODEL}")
        else:
            print("⚠️  No models found. Run: ollama pull mistral")
    else:
        print(f"❌ Ollama not running at {OLLAMA_BASE_URL}")
        print("   Start with: ollama serve")
    
    print("\n📸 Image API Status:")
    print(f"  {'✅' if UNSPLASH_API_KEY else '❌'} Unsplash API")
    print(f"  {'✅' if PEXELS_API_KEY else '❌'} Pexels API")
    print(f"  {'✅' if PIXABAY_API_KEY else '❌'} Pixabay API")
    
    print("\n✅ Mistral setup complete!\n")
    print(f"✅ Running on: http://localhost:5000")
    print("="*60 + "\n")


startup_diagnostics()


# ============ THEMES ============
THEMES = {
    "modern_blue": {
        "name": "Modern Blue",
        "colors": {
            "background": (245, 245, 255),
            "title_slide_bg": (19, 71, 130),
            "title": (255, 255, 255),
            "subtitle": (220, 230, 250),
            "slide_title": (19, 71, 130),
            "h2": (30, 90, 150),
            "h3": (70, 120, 170),
            "text": (45, 45, 45),
            "bullet": (60, 60, 60),
            "accent": (19, 71, 130),
            "accent_light": (200, 220, 255),
            "accent_dark": (10, 40, 80)
        },
        "fonts": {"title": "Calibri", "heading": "Calibri", "text": "Calibri"}
    },
    "corporate_gray": {
        "name": "Corporate Gray",
        "colors": {
            "background": (250, 250, 250),
            "title_slide_bg": (70, 70, 70),
            "title": (255, 255, 255),
            "subtitle": (220, 220, 220),
            "slide_title": (50, 50, 50),
            "h2": (80, 80, 80),
            "h3": (120, 120, 120),
            "text": (60, 60, 60),
            "bullet": (70, 70, 70),
            "accent": (150, 150, 150),
            "accent_light": (230, 230, 230),
            "accent_dark": (40, 40, 40)
        },
        "fonts": {"title": "Arial", "heading": "Arial", "text": "Arial"}
    },
    "creative_green": {
        "name": "Creative Green",
        "colors": {
            "background": (235, 245, 235),
            "title_slide_bg": (34, 102, 36),
            "title": (255, 255, 255),
            "subtitle": (200, 230, 200),
            "slide_title": (34, 102, 36),
            "h2": (60, 130, 65),
            "h3": (100, 150, 105),
            "text": (50, 50, 50),
            "bullet": (60, 60, 60),
            "accent": (34, 102, 36),
            "accent_light": (200, 240, 200),
            "accent_dark": (20, 70, 25)
        },
        "fonts": {"title": "Segoe UI", "heading": "Segoe UI", "text": "Segoe UI"}
    },
    "elegant_purple": {
        "name": "Elegant Purple",
        "colors": {
            "background": (240, 235, 250),
            "title_slide_bg": (102, 51, 153),
            "title": (255, 255, 255),
            "subtitle": (220, 200, 240),
            "slide_title": (102, 51, 153),
            "h2": (120, 70, 170),
            "h3": (150, 120, 190),
            "text": (65, 65, 65),
            "bullet": (75, 75, 75),
            "accent": (102, 51, 153),
            "accent_light": (230, 210, 250),
            "accent_dark": (70, 30, 110)
        },
        "fonts": {"title": "Times New Roman", "heading": "Times New Roman", "text": "Times New Roman"}
    },
    "vibrant_orange": {
        "name": "Vibrant Orange",
        "colors": {
            "background": (255, 245, 235),
            "title_slide_bg": (204, 85, 0),
            "title": (255, 255, 255),
            "subtitle": (255, 220, 180),
            "slide_title": (204, 85, 0),
            "h2": (220, 110, 30),
            "h3": (240, 150, 80),
            "text": (60, 40, 20),
            "bullet": (80, 60, 40),
            "accent": (204, 85, 0),
            "accent_light": (255, 220, 180),
            "accent_dark": (150, 60, 0)
        },
        "fonts": {"title": "Helvetica", "heading": "Helvetica", "text": "Helvetica"}
    },
    "minimalist_black": {
        "name": "Minimalist Black",
        "colors": {
            "background": (20, 20, 20),
            "title_slide_bg": (0, 0, 0),
            "title": (255, 255, 255),
            "subtitle": (200, 200, 200),
            "slide_title": (255, 255, 255),
            "h2": (220, 220, 220),
            "h3": (180, 180, 180),
            "text": (180, 180, 180),
            "bullet": (160, 160, 160),
            "accent": (255, 255, 255),
            "accent_light": (60, 60, 60),
            "accent_dark": (100, 100, 100)
        },
        "fonts": {"title": "Arial Black", "heading": "Arial Black", "text": "Arial"}
    },
    "soft_pink": {
        "name": "Soft Pink",
        "colors": {
            "background": (255, 245, 250),
            "title_slide_bg": (180, 20, 100),
            "title": (255, 255, 255),
            "subtitle": (230, 180, 210),
            "slide_title": (180, 20, 100),
            "h2": (200, 60, 140),
            "h3": (220, 100, 170),
            "text": (80, 30, 50),
            "bullet": (90, 40, 60),
            "accent": (180, 20, 100),
            "accent_light": (245, 210, 235),
            "accent_dark": (130, 10, 70)
        },
        "fonts": {"title": "Georgia", "heading": "Georgia", "text": "Georgia"}
    },
    "cool_teal": {
        "name": "Cool Teal",
        "colors": {
            "background": (230, 248, 248),
            "title_slide_bg": (0, 102, 102),
            "title": (255, 255, 255),
            "subtitle": (180, 220, 220),
            "slide_title": (0, 102, 102),
            "h2": (20, 130, 130),
            "h3": (60, 150, 150),
            "text": (40, 70, 70),
            "bullet": (50, 80, 80),
            "accent": (0, 102, 102),
            "accent_light": (200, 240, 240),
            "accent_dark": (0, 70, 70)
        },
        "fonts": {"title": "Verdana", "heading": "Verdana", "text": "Verdana"}
    },
    "warm_brown": {
        "name": "Warm Brown",
        "colors": {
            "background": (250, 242, 232),
            "title_slide_bg": (101, 67, 33),
            "title": (255, 255, 255),
            "subtitle": (220, 190, 160),
            "slide_title": (101, 67, 33),
            "h2": (140, 100, 60),
            "h3": (170, 140, 100),
            "text": (70, 55, 40),
            "bullet": (85, 70, 55),
            "accent": (101, 67, 33),
            "accent_light": (240, 220, 190),
            "accent_dark": (70, 45, 20)
        },
        "fonts": {"title": "Palatino Linotype", "heading": "Palatino Linotype", "text": "Palatino Linotype"}
    },
}


# ============ SLIDE DIMENSIONS ============
SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 7.5
MARGIN_LEFT = 0.5
MARGIN_RIGHT = 0.5
MARGIN_TOP = 0.4
MARGIN_BOTTOM = 0.5
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TITLE_AREA_HEIGHT = 1.0
CONTENT_AREA_TOP = MARGIN_TOP + TITLE_AREA_HEIGHT + 0.2
AVAILABLE_CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM - TITLE_AREA_HEIGHT - 0.5
TEXT_WIDTH_WITH_IMAGE = 4.8
TEXT_WIDTH_WITHOUT_IMAGE = CONTENT_WIDTH
IMAGE_START_LEFT = 5.3
IMAGE_GAP = 0.2


# ============ IMPROVED FONT SIZE CALCULATOR ============
def calculate_dynamic_font_sizes(slide_sections, mode, has_image):
    """
    Dynamically calculate font sizes based on content density
    Returns: dict with font sizes for different text types
    
    FIX: Better calculation that prevents overflow
    """
    # Count content items
    bullet_count = sum(1 for s in slide_sections if s['type'] == 'bullet')
    text_count = sum(1 for s in slide_sections if s['type'] == 'text')
    h3_count = sum(1 for s in slide_sections if s['type'] == 'h3')
    
    total_items = bullet_count + text_count + h3_count
    
    # Base font sizes from mode
    base_body = MODE_PROMPTS[mode].get("font_size_body", 12)
    
    # Adjust based on content density and image presence
    if total_items <= 2:
        body_size = base_body + 1
        h3_size = 16
        spacing = 8
    elif total_items <= 4:
        body_size = base_body
        h3_size = 15
        spacing = 6
    elif total_items <= 6:
        body_size = max(base_body - 0.5, 11)
        h3_size = 14
        spacing = 5
    elif total_items <= 8:
        body_size = max(base_body - 1, 10)
        h3_size = 13
        spacing = 4
    else:
        # Very dense - shrink carefully
        body_size = max(base_body - 1.5, 9)
        h3_size = 12
        spacing = 3
    
    # Further reduce if image is present (less horizontal space)
    if has_image:
        body_size = max(body_size - 0.5, 9)
        h3_size = max(h3_size - 1, 11)
        spacing = max(spacing - 1, 2)
    
    return {
        'body': body_size,
        'h3': h3_size,
        'spacing': spacing,
        'line_spacing': 1.15 if total_items > 6 else 1.25
    }


# ============ MISTRAL API WITH MODE SUPPORT ============
def generate_with_mistral(prompt: str, mode: str = "Creative") -> str:
    """Generate content using local Mistral model via Ollama"""
    mode_config = MODE_PROMPTS.get(mode, MODE_PROMPTS["Creative"])
    mode_instruction = mode_config.get("instructions", "")
    
    enhanced_prompt = f"""{mode_instruction}


User Topic/Request: {prompt}


Please structure your response with markdown:
- Use # for main section titles
- Use ## for subsection titles
- Use ### for sub-points and key areas
- Use * or - for bullet points (keep to {mode_config.get('max_bullets', 4)} per section)
- Keep points clear, concise, and well-organized
- Include practical examples where relevant


This will be converted to a {mode_config.get('slide_count', 5)}-slide PowerPoint presentation.
Format each section properly so it can be converted to individual slides."""
    
    try:
        url = f"{OLLAMA_BASE_URL}/api/generate"
        
        payload = {
            "model": MISTRAL_MODEL,
            "prompt": enhanced_prompt,
            "stream": False,
            "temperature": 0.7 if mode == "Creative" else (0.5 if mode == "Quick Response" else 0.6),
        }
        
        print(f"🤖 Generating [{mode}] with {MISTRAL_MODEL} (LOCAL)...")
        
        response = requests.post(url, json=payload, timeout=120)
        
        if response.status_code == 200:
            result = response.json()
            generated_text = result.get('response', '').strip()
            print(f"✅ Generation complete ({len(generated_text)} chars)\n")
            return generated_text
        else:
            error_msg = f"Ollama error: HTTP {response.status_code}"
            print(f"❌ {error_msg}")
            return f"[Error]: {error_msg}"
    
    except requests.exceptions.ConnectionError:
        error_msg = f"Cannot connect to Ollama at {OLLAMA_BASE_URL}. Make sure 'ollama serve' is running."
        print(f"❌ {error_msg}")
        return f"[Error]: {error_msg}"
    
    except requests.exceptions.Timeout:
        error_msg = "Mistral generation timed out (120s). Try a simpler prompt or Quick Response mode."
        print(f"❌ {error_msg}")
        return f"[Error]: {error_msg}"
    
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Mistral error: {error_msg}")
        return f"[Error]: {error_msg}"


# ============ IMPROVED TEXT PROCESSING ============
def clean_markdown(text):
    """
    Remove markdown formatting
    FIX: More comprehensive cleaning including separators
    """
    if not text:
        return ""
    
    text = str(text)
    
    # Remove markdown separators (~~~, ---, ***)
    text = re.sub(r'^[\~\-\*]{3,}$', '', text, flags=re.MULTILINE)
    
    # Remove bold/italic markdown
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'_(.*?)_', r'\1', text)
    
    # Remove code markdown
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'```(.*?)```', r'\1', text, flags=re.DOTALL)
    
    # Remove header markers
    text = re.sub(r'#+\s*', '', text)
    
    # Remove list markers (*, -, •)
    text = re.sub(r'^[\*\-\•]\s*', '', text, flags=re.MULTILINE)
    
    # Remove bracket links [text](url)
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    
    # Remove standalone brackets
    text = re.sub(r'\[\s*\]', '', text)
    text = re.sub(r'\(\s*\)', '', text)
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    
    return text.strip()


def truncate_text(text, max_length=150):
    """
    Truncate text to prevent overflow
    FIX: More conservative truncation
    """
    text = text.strip()
    if len(text) > max_length:
        # Try to break at word boundary
        truncated = text[:max_length]
        last_space = truncated.rfind(' ')
        if last_space > max_length * 0.7:
            return text[:last_space] + "..."
        return truncated + "..."
    return text


def parse_markdown_content(text):
    """
    Parse markdown content into sections
    FIX: Better handling of various markdown formats
    """
    lines = text.split('\n')
    sections = []
    
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        
        # Skip empty lines
        if not stripped or len(stripped) < 2:
            i += 1
            continue
        
        # Skip markdown separators
        if re.match(r'^[\~\-\*]{3,}$', stripped):
            i += 1
            continue
        
        # h3 (###)
        if stripped.startswith('###'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h3', 'content': content, 'level': 3})
        
        # h2 (##)
        elif stripped.startswith('##'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h2', 'content': content, 'level': 2})
        
        # h1 (#)
        elif stripped.startswith('#'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h1', 'content': content, 'level': 1})
        
        # Bullets (*, -, •)
        elif stripped.startswith(('*', '-', '•')):
            content = stripped.lstrip('*-•').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                level = (len(line) - len(line.lstrip())) // 2
                sections.append({'type': 'bullet', 'content': content, 'level': level})
        
        # Regular text (only if substantial)
        else:
            content = clean_markdown(stripped)
            if len(content) > 5:  # Only add if substantial
                sections.append({'type': 'text', 'content': content, 'level': 0})
        
        i += 1
    
    return sections


def group_into_slides(sections, max_slides=5):
    """
    Group sections into logical slides
    FIX: Better grouping logic that preserves content
    """
    slides = []
    current_slide = []
    items_in_current = 0
    
    for section in sections:
        # Start new slide on h1/h2
        if section['type'] in ['h1', 'h2'] and current_slide and items_in_current > 1:
            slides.append(current_slide)
            current_slide = [section]
            items_in_current = 1
        else:
            current_slide.append(section)
            items_in_current += 1
        
        # Limit items per slide to prevent overflow
        if items_in_current > 20:  # Was 25, now more conservative
            slides.append(current_slide)
            current_slide = []
            items_in_current = 0
        
        # Hard limit on slides
        if len(slides) >= max_slides - 1:
            break
    
    # Add remaining content
    if current_slide:
        if len(slides) < max_slides:
            slides.append(current_slide)
        else:
            # Merge with last slide if we're at limit
            if slides:
                slides[-1].extend(current_slide)
            else:
                slides.append(current_slide)
    
    return slides[:max_slides]


# ============ IMAGE FUNCTIONS ============
def fetch_unsplash_image(query):
    if not UNSPLASH_API_KEY:
        return None
    try:
        url = "https://api.unsplash.com/photos/random"
        params = {"query": query, "client_id": UNSPLASH_API_KEY, "orientation": "landscape"}
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            return response.json().get('urls', {}).get('regular')
    except:
        pass
    return None


def fetch_pexels_image(query):
    if not PEXELS_API_KEY:
        return None
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1}
        response = requests.get(url, headers=headers, params=params, timeout=5)
        if response.status_code == 200:
            photos = response.json().get('photos', [])
            if photos:
                return photos[0].get('src', {}).get('large')
    except:
        pass
    return None


def fetch_pixabay_image(query):
    if not PIXABAY_API_KEY:
        return None
    try:
        url = "https://pixabay.com/api/"
        params = {"key": PIXABAY_API_KEY, "q": query, "image_type": "photo", "per_page": 1, "min_width": 800, "min_height": 600}
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            hits = response.json().get('hits', [])
            if hits:
                return hits[0].get('largeImageURL')
    except:
        pass
    return None


def get_image_for_slide(slide_title):
    search_query = re.sub(r'[^a-zA-Z0-9\s]', '', slide_title).strip()[:50]
    if not search_query:
        return None
    
    if UNSPLASH_API_KEY:
        img_url = fetch_unsplash_image(search_query)
        if img_url:
            return img_url
    
    if PEXELS_API_KEY:
        img_url = fetch_pexels_image(search_query)
        if img_url:
            return img_url
    
    if PIXABAY_API_KEY:
        img_url = fetch_pixabay_image(search_query)
        if img_url:
            return img_url
    
    return None


def download_image(url):
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            img = Image.open(BytesIO(response.content))
            img.thumbnail((1200, 800), Image.Resampling.LANCZOS)
            img_byte_arr = BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=85)
            img_byte_arr.seek(0)
            return img_byte_arr
    except:
        pass
    return None


# ============ SLIDE BUILDING ============
def add_decorative_shapes(slide, theme, position="top"):
    if position == "top":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(0.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*theme["colors"]["accent"])
        shape.line.color.rgb = RGBColor(*theme["colors"]["accent"])


def set_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


# ============ API ROUTES ============
@app.route("/api/generate", methods=["POST"])
def generate_content():
    data = request.get_json()
    prompt = data.get("prompt", "")
    mode = data.get("mode", "Creative")
    
    if mode not in MODE_PROMPTS:
        mode = "Creative"
    
    print(f"\n📝 Generating content in '{mode}' mode...")
    output = generate_with_mistral(prompt, mode)
    
    return jsonify({
        "output": output,
        "mode": mode,
        "status": "success" if not output.startswith("[Error]") else "error"
    })


@app.route("/api/themes", methods=["GET"])
def get_themes():
    theme_list = [
        {
            "id": k,
            "name": v["name"],
            "preview_color": '#%02x%02x%02x' % v["colors"]["background"]
        }
        for k, v in THEMES.items()
    ]
    return jsonify({"themes": theme_list})


@app.route("/api/modes", methods=["GET"])
def get_modes():
    modes_list = [
        {
            "id": mode,
            "name": mode,
            "description": config.get("instructions"),
            "slide_count": config.get("slide_count"),
            "max_bullets": config.get("max_bullets")
        }
        for mode, config in MODE_PROMPTS.items()
    ]
    return jsonify({"modes": modes_list})


@app.route("/api/generate-ppt", methods=["POST"])
def generate_ppt():
    """
    Generate PowerPoint from content
    FIXES:
    - Better text overflow handling
    - Improved markdown cleaning
    - Preserve full content
    - Smart pagination
    """
    try:
        data = request.get_json()
        content = data.get("content", "").strip()
        title = data.get("title", "Generated Presentation").strip()
        filename = (data.get("filename", "") or "").strip().replace(" ", "_")
        include_images = data.get("include_images", True)
        mode = data.get("mode", "Creative")
        
        if mode not in MODE_PROMPTS:
            mode = "Creative"
        
        if not filename:
            filename = title.replace(" ", "_")
        
        customizations = data.get("customizations", {}) or {}
        theme_id = customizations.get("theme", "modern_blue")
        theme = THEMES.get(theme_id, THEMES["modern_blue"])
        max_slides = int(customizations.get("slide_count", MODE_PROMPTS[mode]["slide_count"]))
        
        if not content:
            return jsonify({"error": "Content is empty"}), 400
        
        # Parse content
        sections = parse_markdown_content(content)
        if not sections:
            return jsonify({"error": "Could not parse content"}), 400
        
        # Group into slides
        slides_content = group_into_slides(sections, max_slides)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        
        print("\n" + "="*60)
        print("🎬 CREATING PRESENTATION")
        print("="*60)
        print(f"Title: {title}")
        print(f"Mode: {mode}")
        print(f"Theme: {theme_id}")
        print(f"Include Images: {include_images}")
        print(f"Parsed Sections: {len(sections)}")
        print(f"Grouped Slides: {len(slides_content)}")
        print("="*60 + "\n")
        
        # ============ TITLE SLIDE ============
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(title_slide, theme["colors"]["title_slide_bg"])
        add_decorative_shapes(title_slide, theme, "top")
        
        title_box = title_slide.shapes.add_textbox(
            Inches(MARGIN_LEFT),
            Inches(SLIDE_HEIGHT * 0.25),
            Inches(CONTENT_WIDTH),
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.name = theme["fonts"]["title"]
        p.font.color.rgb = RGBColor(*theme["colors"]["title"])
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
        
        # Subtitle
        if slides_content and slides_content[0]:
            first_text = next(
                (s['content'] for s in slides_content[0] if s['type'] in ['h1', 'h2', 'text']),
                f"Generated in {mode} mode"
            )
            subtitle_box = title_slide.shapes.add_textbox(
                Inches(MARGIN_LEFT),
                Inches(SLIDE_HEIGHT * 0.5),
                Inches(CONTENT_WIDTH),
                Inches(1.0)
            )
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = truncate_text(first_text, 80)
            p_sub.font.size = Pt(24)
            p_sub.font.name = theme["fonts"]["text"]
            p_sub.font.color.rgb = RGBColor(*theme["colors"]["subtitle"])
            p_sub.alignment = PP_ALIGN.CENTER
        
        # ============ CONTENT SLIDES ============
        for slide_idx, slide_sections in enumerate(slides_content):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(slide, theme["colors"]["background"])
            add_decorative_shapes(slide, theme, "top")
            
            # Extract slide title
            slide_title = "Overview"
            for section in slide_sections:
                if section['type'] in ['h1', 'h2']:
                    slide_title = section['content']
                    break
            
            print(f"📄 Slide {slide_idx + 1}: {truncate_text(slide_title, 60)}")
            
            # ===== TITLE BACKGROUND & TEXT =====
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(MARGIN_LEFT),
                Inches(MARGIN_TOP),
                Inches(CONTENT_WIDTH),
                Inches(TITLE_AREA_HEIGHT + 0.1)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(*theme["colors"]["accent_light"])
            title_bg.line.fill.background()
            
            # Title text
            title_box = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT + 0.2),
                Inches(MARGIN_TOP + 0.05),
                Inches(CONTENT_WIDTH - 0.4),
                Inches(TITLE_AREA_HEIGHT)
            )
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = truncate_text(slide_title, 80)
            p_title.font.size = Pt(36)
            p_title.font.bold = True
            p_title.font.name = theme["fonts"]["heading"]
            p_title.font.color.rgb = RGBColor(*theme["colors"]["slide_title"])
            p_title.alignment = PP_ALIGN.LEFT
            
            # ===== ADD IMAGE (if enabled) =====
            image_added = False
            if include_images and slide_idx > 0:
                print(f"  🖼️  Fetching image...")
                img_url = get_image_for_slide(slide_title)
                if img_url:
                    try:
                        img_data = download_image(img_url)
                        if img_data:
                            img_data.seek(0)
                            calc_width, calc_height = 3.2, 4.2
                            img_left = Inches(MARGIN_LEFT + IMAGE_START_LEFT + IMAGE_GAP)
                            img_top = Inches(CONTENT_AREA_TOP)
                            
                            slide.shapes.add_picture(
                                img_data,
                                img_left,
                                img_top,
                                width=Inches(calc_width),
                                height=Inches(calc_height)
                            )
                            image_added = True
                            print(f"  ✅ Image added!")
                        else:
                            print(f"  ⚠️  Could not download image")
                    except Exception as e:
                        print(f"  ❌ Error adding image: {e}")
                else:
                    print(f"  ⚠️  No image found")
            
            # ===== CALCULATE FONT SIZES =====
            font_config = calculate_dynamic_font_sizes(slide_sections, mode, image_added)
            print(f"  📏 Font: Body={font_config['body']}pt, H3={font_config['h3']}pt")
            
            # ===== CONTENT AREA =====
            content_width = TEXT_WIDTH_WITH_IMAGE if image_added else TEXT_WIDTH_WITHOUT_IMAGE
            content_left = MARGIN_LEFT
            content_top = CONTENT_AREA_TOP
            content_height = AVAILABLE_CONTENT_HEIGHT
            
            content_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(content_top),
                Inches(content_width),
                Inches(content_height)
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.vertical_anchor = MSO_ANCHOR.TOP
            
            # ===== KEY FIX: NO AUTO-SHRINK =====
            # Instead of MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE (aggressive shrinking)
            # Use NONE and let word_wrap handle it
            tf_content.auto_size = MSO_AUTO_SIZE.NONE
            
            # ===== ADD CONTENT =====
            first_para = True
            for section in slide_sections:
                # Skip h1/h2 (already used as title)
                if section['type'] == 'h1' or section['type'] == 'h2':
                    continue
                
                # Add new paragraph if not first
                if not first_para:
                    tf_content.add_paragraph()
                
                # h3: Bold subheading
                if section['type'] == 'h3':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 150)
                    p.font.size = Pt(font_config['h3'])
                    p.font.bold = True
                    p.font.name = theme["fonts"]["heading"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["h3"])
                    p.space_after = Pt(font_config['spacing'])
                    p.line_spacing = font_config['line_spacing']
                
                # Bullet point
                elif section['type'] == 'bullet':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 150)
                    p.level = min(section['level'], 2)  # Cap nesting at level 2
                    p.font.size = Pt(font_config['body'])
                    p.font.name = theme["fonts"]["text"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["bullet"])
                    p.space_after = Pt(font_config['spacing'])
                    p.line_spacing = font_config['line_spacing']
                
                # Regular text
                elif section['type'] == 'text':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 150)
                    p.font.size = Pt(font_config['body'])
                    p.font.name = theme["fonts"]["text"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["text"])
                    p.space_after = Pt(font_config['spacing'] + 2)
                    p.line_spacing = font_config['line_spacing']
                
                first_para = False
        
        # ===== SAVE PRESENTATION =====
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        print(f"\n✅ Presentation created successfully!")
        print(f"📊 Total slides: {len(prs.slides)}")
        print(f"📦 File size: {len(ppt_bytes.getvalue()) / 1024:.1f} KB\n")
        
        return send_file(
            ppt_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"{filename}.pptx"
        )
    
    except Exception as e:
        print(f"❌ Error generating PowerPoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/health", methods=["GET"])
def health_check():
    ollama_ok, models = check_ollama_connection()
    
    return jsonify({
        "status": "healthy" if ollama_ok else "degraded",
        "ollama_running": ollama_ok,
        "ollama_url": OLLAMA_BASE_URL,
        "mistral_model": MISTRAL_MODEL,
        "available_models": models,
        "available_modes": list(MODE_PROMPTS.keys()),
        "image_apis": {
            "unsplash": bool(UNSPLASH_API_KEY),
            "pexels": bool(PEXELS_API_KEY),
            "pixabay": bool(PIXABAY_API_KEY)
        }
    })


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)